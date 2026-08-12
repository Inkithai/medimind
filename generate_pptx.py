#!/usr/bin/env python3
"""
Generates PRESENTATION.pptx — a professional 16:9 widescreen PowerPoint presentation
for MediMind with custom dark medical styling, structured cards, diagrams, tables,
and comprehensive Speaker Notes for judges.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------------------
# COLOR PALETTE (Dark Medical Theme)
# ---------------------------------------------------------------------------
BG_DARK = RGBColor(15, 23, 42)       # #0F172A Deep Navy / Slate background
CARD_BG = RGBColor(30, 41, 59)       # #1E293B Card background
CARD_BORDER = RGBColor(51, 65, 85)   # #334155 Subtle border
TEXT_WHITE = RGBColor(248, 250, 252) # #F8FAFC Primary title text
TEXT_MUTED = RGBColor(148, 163, 184) # #94A3B8 Secondary / body text
CYAN_ACCENT = RGBColor(6, 182, 212)  # #06B6D4 Teal / Cyan highlight
WARNING_AMBER = RGBColor(245, 158, 11) # #F59E0B Warning accent
SUCCESS_EMERALD = RGBColor(16, 185, 129) # #10B981 Success accent


def set_slide_background(slide):
    """Sets a solid deep navy/slate background on the slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK


def add_header(slide, slide_num: int, category: str, title: str, subtitle: str):
    """Adds a standardized professional slide header zone."""
    # Top Category / Number label
    cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(11.7), Inches(0.35))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = f"SLIDE {slide_num} OF 8  •  {category.upper()}"
    p_cat.font.name = "Arial"
    p_cat.font.size = Pt(11)
    p_cat.font.bold = True
    p_cat.font.color.rgb = CYAN_ACCENT

    # Main Slide Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.80), Inches(11.7), Inches(0.70))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title
    p_title.font.name = "Arial"
    p_title.font.size = Pt(26)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_WHITE

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.50), Inches(11.7), Inches(0.45))
    tf_sub = sub_box.text_frame
    tf_sub.word_wrap = True
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = subtitle
    p_sub.font.name = "Arial"
    p_sub.font.size = Pt(14)
    p_sub.font.color.rgb = TEXT_MUTED

    # Header Bottom Divider Line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), Inches(2.05), Inches(11.733), Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = CARD_BORDER
    line.line.fill.background()


def add_card(slide, left, top, width, height, tag: str, title: str, body_text: str, accent_color=CYAN_ACCENT):
    """Adds a visually styled card with tag badge, bold title, and body text."""
    # Background shape
    card_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card_shape.fill.solid()
    card_shape.fill.fore_color.rgb = CARD_BG
    card_shape.line.color.rgb = CARD_BORDER
    card_shape.line.width = Pt(1.5)

    # Accent left border stripe
    stripe = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(0.08), height)
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent_color
    stripe.line.fill.background()

    # Text box inside card
    txbox = slide.shapes.add_textbox(
        left + Inches(0.2),
        top + Inches(0.15),
        width - Inches(0.35),
        height - Inches(0.3)
    )
    tf = txbox.text_frame
    tf.word_wrap = True

    # Tag Badge paragraph
    p_tag = tf.paragraphs[0]
    p_tag.text = tag.upper()
    p_tag.font.name = "Arial"
    p_tag.font.size = Pt(10)
    p_tag.font.bold = True
    p_tag.font.color.rgb = accent_color
    p_tag.space_after = Pt(4)

    # Title paragraph
    p_title = tf.add_paragraph()
    p_title.text = title
    p_title.font.name = "Arial"
    p_title.font.size = Pt(15)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_WHITE
    p_title.space_after = Pt(8)

    # Body paragraph
    p_body = tf.add_paragraph()
    p_body.text = body_text
    p_body.font.name = "Arial"
    p_body.font.size = Pt(12)
    p_body.font.color.rgb = TEXT_MUTED
    p_body.line_spacing = 1.25


def add_speaker_note(slide, note_text: str):
    """Adds speaker notes to the presentation slide."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = note_text


def build_presentation():
    prs = Presentation()
    # 16:9 Widescreen dimensions (13.333" x 7.5")
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # =========================================================================
    # SLIDE 1: PROBLEM
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1)
    add_header(
        slide1,
        1, "Real-World Challenge",
        "The Broken Medical Record Experience",
        "Why managing personal health records is dangerous, fragmented, and overwhelming"
    )

    # 2x2 grid cards
    add_card(
        slide1,
        Inches(0.8), Inches(2.35), Inches(5.6), Inches(2.1),
        "Problem 01", "Fragmented Health Silos",
        "Patients receive unstructured records from different hospitals, labs, and specialists—paper prescriptions, PDF blood tests, discharge notes, and scanned images. Critical medical history is trapped in disconnected silos.",
        WARNING_AMBER
    )
    add_card(
        slide1,
        Inches(6.8), Inches(2.35), Inches(5.7), Inches(2.1),
        "Problem 02", "Clinical Jargon & Confusion",
        "Medical documents are written for clinicians, not patients. Complex terminology, Latin dosage abbreviations, and obscure reference ranges leave patients and caregivers anxious and unsure about their own health status.",
        WARNING_AMBER
    )
    add_card(
        slide1,
        Inches(0.8), Inches(4.75), Inches(5.6), Inches(2.1),
        "Problem 03", "Polypharmacy & Medication Errors",
        "When visiting multiple specialists, patients rarely bring a comprehensive drug list. Overlooked drug-drug interactions, contraindications against allergies, and duplicate therapies cause preventable adverse drug events (ADEs).",
        WARNING_AMBER
    )
    add_card(
        slide1,
        Inches(6.8), Inches(4.75), Inches(5.7), Inches(2.1),
        "Problem 04", "Lost Longitudinal Biomarkers",
        "A single lab report tells only half the story. Tracking vital biomarker trajectories over years (e.g., eGFR, HbA1c, cholesterol, liver enzymes) is manual, tedious, and frequently ignored until acute symptoms appear.",
        WARNING_AMBER
    )

    add_speaker_note(
        slide1,
        "Every one of us has experienced—or helped a family member navigate—the chaos of scattered medical records. When health data is locked inside unreadable scans and disconnected PDFs, patients face real safety risks from overlooked drug interactions and missed lab trends. MediMind solves this by turning passive medical documents into an intelligent, active health concierge."
    )

    # =========================================================================
    # SLIDE 2: SOLUTION
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2)
    add_header(
        slide2,
        2, "Solution Overview",
        "Solution — MediMind AI Concierge",
        "Transforming chaotic medical documents into a structured, proactive health timeline"
    )

    add_card(
        slide2,
        Inches(0.8), Inches(2.35), Inches(5.6), Inches(2.1),
        "Value 01", "Intelligent Medical Concierge",
        "An AI-powered, privacy-first platform that ingests unstructured medical documents and synthesizes them into an actionable health dashboard—bridging clinical documentation and patient understanding without losing medical precision.",
        SUCCESS_EMERALD
    )
    add_card(
        slide2,
        Inches(6.8), Inches(2.35), Inches(5.7), Inches(2.1),
        "Value 02", "Automated Structured Extraction",
        "Upload any combination of prescriptions, lab reports, clinical notes, or discharge summaries. MediMind automatically categorizes documents and extracts structured entities (medications, lab biomarkers, diagnoses, and allergies).",
        SUCCESS_EMERALD
    )
    add_card(
        slide2,
        Inches(0.8), Inches(4.75), Inches(5.6), Inches(2.1),
        "Value 03", "Unified Longitudinal Timeline",
        "Synthesizes isolated visits into a chronological health history. Every prescription, lab test, and physician encounter is linked into a unified patient timeline with full traceability back to original source documents.",
        SUCCESS_EMERALD
    )
    add_card(
        slide2,
        Inches(6.8), Inches(4.75), Inches(5.7), Inches(2.1),
        "Value 04", "Proactive Safety Cross-Checks",
        "Continuously audits the patient's active medication list against known allergies, chronic conditions, and drug-drug interactions, while tracking lab biomarker trajectories to flag clinical abnormalities.",
        SUCCESS_EMERALD
    )

    add_speaker_note(
        slide2,
        "MediMind is not just an OCR scanner or a general chatbot. It is a purpose-built medical intelligence platform. When a user uploads a stack of medical records, MediMind extracts the clinical facts, constructs a coherent timeline, cross-checks every medication for safety, and allows patients to ask questions against their own health history."
    )

    # =========================================================================
    # SLIDE 3: AI CORE
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3)
    add_header(
        slide3,
        3, "Deep Clinical Intelligence",
        "AI Core & RAG Pipeline",
        "Multimodal OCR, structured JSON recovery ladders, and zero-network ONNX embeddings"
    )

    # 3 Top Columns (Width ~3.7 each)
    add_card(
        slide3,
        Inches(0.8), Inches(2.35), Inches(3.7), Inches(2.2),
        "1. OCR & Vision", "Multimodal Understanding",
        "Combines native PDF parsing (PyMuPDF / pdfplumber) with multimodal vision models (qwen3.6-27b, gemini-3.6-flash) to decode handwritten notes, prescription photos, and hospital scans.",
        CYAN_ACCENT
    )
    add_card(
        slide3,
        Inches(4.8), Inches(2.35), Inches(3.7), Inches(2.2),
        "2. LLM Extraction", "Resilient JSON Ladder",
        "Multi-rung structured decoding ladder (_completion_resilient). Uses constrained JSON Schema with automated <think> reasoning-tag suppression probes and an OpenRouter auto-fallback circuit breaker.",
        CYAN_ACCENT
    )
    add_card(
        slide3,
        Inches(8.8), Inches(2.35), Inches(3.7), Inches(2.2),
        "3. Local RAG", "Zero-Network Embeddings",
        "Runs all-MiniLM-L6-v2 locally in-process via ONNX Runtime—no external embedding network calls required. Embeds structured timeline chunks in ChromaDB for hallucination-free Q&A.",
        CYAN_ACCENT
    )

    # 2 Bottom Columns
    add_card(
        slide3,
        Inches(0.8), Inches(4.75), Inches(5.6), Inches(2.1),
        "4. Clinical Extraction", "Structured Health Entity Schema",
        "Constrained Pydantic schemas normalize every document into strict JSON: active medications with dosages and timing, lab biomarker panels with reference units, clinical diagnoses, and verified allergy lists.",
        CYAN_ACCENT
    )
    add_card(
        slide3,
        Inches(6.8), Inches(4.75), Inches(5.7), Inches(2.1),
        "5. Safety Engine", "Automated Clinical Safety Audit",
        "Runs automated polypharmacy audits (cross_check_prescriptions) coded by severity (High / Moderate / Precaution) and tracks longitudinal lab biomarker percentage shifts (track_lab_trends).",
        CYAN_ACCENT
    )

    add_speaker_note(
        slide3,
        "Our AI core is engineered for clinical resilience. Medical documents are messy, and standard LLMs often break on formatting or leak verbose reasoning tags. We built a multi-rung structured recovery ladder with automated reasoning-tag suppression, local zero-network ONNX embeddings for RAG, and an automated safety engine that audits prescriptions against patient allergies and lab trends."
    )

    # =========================================================================
    # SLIDE 4: KEY FEATURES
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4)
    add_header(
        slide4,
        4, "Product Capabilities",
        "Key Features & Comprehensive Suite",
        "A complete suite empowering patients and caregivers to take control of their medical history"
    )

    # 6 feature cards (2 rows x 3 columns)
    add_card(
        slide4,
        Inches(0.8), Inches(2.35), Inches(3.7), Inches(2.15),
        "Feature 1", "Async Document Upload",
        "Drag-and-drop batch upload for PDFs, PNGs, and JPGs. Non-blocking background worker pool (/api/v1/documents?async=true) with real-time job progress tracking.",
        CYAN_ACCENT
    )
    add_card(
        slide4,
        Inches(4.8), Inches(2.35), Inches(3.7), Inches(2.15),
        "Feature 2", "Report Extraction",
        "Automatic document classification and structured extraction of medicine dosages, frequencies, lab biomarkers, reference ranges, and doctor notes.",
        CYAN_ACCENT
    )
    add_card(
        slide4,
        Inches(8.8), Inches(2.35), Inches(3.7), Inches(2.15),
        "Feature 3", "Patient Timeline",
        "Interactive chronological health history (TimelineView.tsx) unifying all past visits, prescriptions, and test results with source PDF links.",
        CYAN_ACCENT
    )

    add_card(
        slide4,
        Inches(0.8), Inches(4.75), Inches(3.7), Inches(2.15),
        "Feature 4", "Medicine History",
        "Consolidated medication dashboard separating active prescriptions from historical treatments, showing dosages, start/end dates, and adherence context.",
        CYAN_ACCENT
    )
    add_card(
        slide4,
        Inches(4.8), Inches(4.75), Inches(3.7), Inches(2.15),
        "Feature 5", "Safety Cross-Checks",
        "Dedicated Cross-Check Safety Report (CrossCheckView.tsx) flagging polypharmacy drug-drug interactions, allergy warnings, and organ precautions.",
        CYAN_ACCENT
    )
    add_card(
        slide4,
        Inches(8.8), Inches(4.75), Inches(3.7), Inches(2.15),
        "Feature 6", "Conversational AI Q&A",
        "Multi-turn conversational concierge (conversation.py) with history query rewriting to answer natural language questions against stored medical facts.",
        CYAN_ACCENT
    )

    add_speaker_note(
        slide4,
        "From a user's perspective, MediMind offers six core features: seamless multi-file batch uploads, automated clinical data extraction, an intuitive interactive timeline, a consolidated medicine history, severity-graded safety cross-checks, and a conversational RAG assistant that understands conversational context."
    )

    # =========================================================================
    # SLIDE 5: SYSTEM ARCHITECTURE
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5)
    add_header(
        slide5,
        5, "Production-Ready Engineering",
        "System Architecture",
        "Modular, scalable, privacy-first full-stack TypeScript & Python architecture"
    )

    add_card(
        slide5,
        Inches(0.8), Inches(2.35), Inches(5.6), Inches(2.1),
        "Frontend SPA", "React 18 + Vite + TypeScript",
        "Responsive SPA built with Tailwind CSS, modular UI components, real-time polling custom hooks (useApi.ts), and clear processing status indicators.",
        CYAN_ACCENT
    )
    add_card(
        slide5,
        Inches(6.8), Inches(2.35), Inches(5.7), Inches(2.1),
        "Backend API", "FastAPI + Uvicorn Async Python",
        "Clean modular services: api.py (routes & auth), medical_extractor.py (LLM & safety ladder), retrieval.py (ONNX RAG), and conversation.py (chat memory).",
        CYAN_ACCENT
    )
    add_card(
        slide5,
        Inches(0.8), Inches(4.75), Inches(5.6), Inches(2.1),
        "Database & Storage", "Supabase Postgres + Cloudinary",
        "Supabase relational tables (patient_snapshots, jobs, chunks) with anonymized token scopes (anon_*) + Cloudinary encrypted cloud storage for source PDFs.",
        CYAN_ACCENT
    )
    add_card(
        slide5,
        Inches(6.8), Inches(4.75), Inches(5.7), Inches(2.1),
        "External AI Services", "Groq + Gemini + OpenRouter Fallback",
        "Universal OpenAI SDK wrapper supporting Groq (gpt-oss-120b, qwen3.6-27b), Google Gemini 3.6-Flash, and automated OpenRouter circuit-breaker fallback.",
        CYAN_ACCENT
    )

    add_speaker_note(
        slide5,
        "Our architecture is clean, modular, and built for production. The React/Vite frontend communicates via REST with a FastAPI Python backend. The AI layer decouples extraction, RAG retrieval, and conversation memory into distinct modules, backed by Supabase Postgres, Cloudinary, and local zero-network ONNX embeddings."
    )

    # =========================================================================
    # SLIDE 6: DEMO / USER FLOW
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6)
    add_header(
        slide6,
        6, "Step-by-Step Journey",
        "Demo / User Flow",
        "How MediMind processes 6 mixed medical files into instant clinical intelligence"
    )

    # Flow Banner Box across top
    flow_box = slide6.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(2.35), Inches(11.7), Inches(1.8)
    )
    flow_box.fill.solid()
    flow_box.fill.fore_color.rgb = RGBColor(9, 14, 26)
    flow_box.line.color.rgb = CYAN_ACCENT
    flow_box.line.width = Pt(1.5)

    tf_flow = flow_box.text_frame
    tf_flow.word_wrap = True
    p_flow1 = tf_flow.paragraphs[0]
    p_flow1.text = "   [ 1. UPLOAD DOCUMENTS ]   ----->   [ 2. OCR / VISION ]   ----->   [ 3. AI RESILIENT LADDER ]   "
    p_flow1.font.name = "Courier New"
    p_flow1.font.size = Pt(14)
    p_flow1.font.bold = True
    p_flow1.font.color.rgb = CYAN_ACCENT
    p_flow1.alignment = PP_ALIGN.CENTER
    p_flow1.space_after = Pt(12)

    p_flow2 = tf_flow.add_paragraph()
    p_flow2.text = "   [ 4. STRUCTURED POSTGRES DATA ]   ----->   [ 5. TIMELINE / MEDS / SAFETY ALERTS / RAG CHAT ]   "
    p_flow2.font.name = "Courier New"
    p_flow2.font.size = Pt(14)
    p_flow2.font.bold = True
    p_flow2.font.color.rgb = TEXT_WHITE
    p_flow2.alignment = PP_ALIGN.CENTER

    # 2 Cards below flow
    add_card(
        slide6,
        Inches(0.8), Inches(4.5), Inches(5.6), Inches(2.35),
        "Step 01-03", "Asynchronous Worker Processing",
        "When a patient uploads 6 files, the server returns immediately with HTTP 202 Accepted. Background worker tasks read, extract, and index documents concurrently without blocking the UI.",
        CYAN_ACCENT
    )
    add_card(
        slide6,
        Inches(6.8), Inches(4.5), Inches(5.7), Inches(2.35),
        "Step 04-05", "Instant Actionable Dashboard",
        "Once complete, the dashboard refreshes automatically to reveal active medications, severity-graded cross-check warnings, interactive lab biomarker graphs, and an AI chat assistant ready for Q&A.",
        CYAN_ACCENT
    )

    add_speaker_note(
        slide6,
        "Let’s walk through the exact user journey. A patient drops a stack of mixed medical documents into MediMind. Our background worker pool ingests and OCRs each file, runs our structured recovery ladder to extract normalized clinical data, checks for drug interactions, and delivers a clean timeline, active medication list, and safety report in seconds."
    )

    # =========================================================================
    # SLIDE 7: INNOVATION & IMPACT
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7)
    add_header(
        slide7,
        7, "Differentiating Value",
        "Innovation & Real-World Impact",
        "Why MediMind stands out in medical AI and who it transforms"
    )

    add_card(
        slide7,
        Inches(0.8), Inches(2.35), Inches(3.7), Inches(4.5),
        "Innovation", "100% Resilient AI Pipeline",
        "Unlike generic AI wrappers that crash on <think> tags or rate limits, MediMind features automated reasoning suppression probes and an auto-activated OpenRouter hard-quota fallback.\n\nRuns zero-network local ONNX embeddings (all-MiniLM-L6-v2) to protect patient privacy.",
        CYAN_ACCENT
    )
    add_card(
        slide7,
        Inches(4.8), Inches(2.35), Inches(3.7), Inches(4.5),
        "Who Benefits", "Patients, Caregivers & Doctors",
        "Empowers elderly and polypharmacy patients managing 5+ daily drugs by consolidating prescriptions from multiple specialists.\n\nSaves clinicians 15+ minutes per consultation by providing a synthesized chronological history and pre-audited drug list.",
        CYAN_ACCENT
    )
    add_card(
        slide7,
        Inches(8.8), Inches(2.35), Inches(3.7), Inches(4.5),
        "Real-World Impact", "Preventing Medication Errors",
        "Directly mitigates adverse drug events (ADEs)—one of the leading causes of preventable hospitalization worldwide.\n\nEmpowers proactive health intervention by highlighting subtle longitudinal lab biomarker percentage shifts before acute symptoms manifest.",
        SUCCESS_EMERALD
    )

    add_speaker_note(
        slide7,
        "What makes MediMind truly innovative is our synthesis of clinical rigor and technical resilience. Unlike generic wrappers, our platform guarantees structured JSON output through reasoning suppression and auto-fallback circuit breakers, while running zero-network local RAG embeddings to protect patient privacy. We are saving doctors time and preventing life-threatening medication errors for patients."
    )

    # =========================================================================
    # SLIDE 8: TECHNOLOGY STACK & FUTURE
    # =========================================================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8)
    add_header(
        slide8,
        8, "Roadmap & Tech Stack",
        "Technology Stack & Future Roadmap",
        "Modern production tooling and strategic expansion into hospital EMR interoperability"
    )

    # Left Box: Tech Stack Table / Overview
    add_card(
        slide8,
        Inches(0.8), Inches(2.35), Inches(5.8), Inches(4.5),
        "Production Tech Stack", "Full-Stack TypeScript & Python 3.11",
        "• FRONTEND: React 18, Vite, TypeScript, Tailwind CSS, Lucide Icons, Axios, modular responsive UI.\n\n"
        "• BACKEND API: Python 3.11, FastAPI, Uvicorn, Pydantic v2, PyMuPDF, pdfplumber, Pillow.\n\n"
        "• LLM & RAG: Groq (gpt-oss-120b/qwen vision), Google Gemini 3.6-Flash, OpenRouter Auto-Fallback, ChromaDB / ONNX MiniLM.\n\n"
        "• STORAGE & CLOUD: Supabase Postgres (anonymized tokens), Cloudinary cloud archiving, Vercel / Render deployment.",
        CYAN_ACCENT
    )

    # Right Box 1: Roadmap EHR
    add_card(
        slide8,
        Inches(6.9), Inches(2.35), Inches(5.6), Inches(2.1),
        "Future Roadmap 01", "EHR/EMR Interoperability (HL7 FHIR)",
        "Expand ingestion beyond PDFs/images to support direct sync with SMART-on-FHIR hospital health records and Apple Health / Google Fit wearable biomarkers.",
        SUCCESS_EMERALD
    )
    # Right Box 2: Roadmap RBAC
    add_card(
        slide8,
        Inches(6.9), Inches(4.75), Inches(5.6), Inches(2.1),
        "Future Roadmap 02", "Clinician Sharing Portals & RBAC",
        "Secure QR-code and link sharing with granular Role-Based Access Control (RBAC) for primary care doctors, specialists, and authorized family members.",
        SUCCESS_EMERALD
    )

    add_speaker_note(
        slide8,
        "Our technology stack is built on modern, type-safe, high-performance tooling—from React and Vite on the frontend to FastAPI, Supabase, and local ONNX embeddings on the backend. As we look ahead, we are expanding MediMind into a universal health OS with HL7 FHIR hospital interoperability, wearable health sync, and secure clinician sharing portals. Thank you!"
    )

    # Save alongside the other pitch/deploy documents in docs/, resolved
    # relative to this script rather than the caller's working directory so
    # the output lands in the same place no matter where it is invoked from.
    output_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "docs")
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, "PRESENTATION.pptx")
    prs.save(output_path)
    print(f"Successfully generated {output_path}")


if __name__ == "__main__":
    build_presentation()
